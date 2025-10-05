"""Utility script to evaluate the bullish ELN structure on selected US equities.

This module assembles assumed market inputs, prices the structured payoff using
Black-Scholes analytics for a two-name equal-weight basket, and exports
supporting artefacts (charts, data tables, Excel workbook, PPT placeholder).

The script is deliberately written to be easy to read and adapt: assumptions are
collected near the top and the functions are self-contained.
"""
from __future__ import annotations

import math
from dataclasses import dataclass
from pathlib import Path
from typing import Dict, List, Tuple

try:
    import matplotlib.pyplot as plt  # type: ignore
except ImportError:  # pragma: no cover - optional dependency
    plt = None  # type: ignore

try:
    import numpy as np  # type: ignore
except ImportError:  # pragma: no cover - optional dependency
    np = None  # type: ignore

try:
    import pandas as pd  # type: ignore
except ImportError:  # pragma: no cover - optional dependency
    pd = None  # type: ignore

try:
    from pptx import Presentation  # type: ignore
    from pptx.util import Inches  # type: ignore
    PPTX_AVAILABLE = True
except ImportError:  # pragma: no cover - optional dependency
    PPTX_AVAILABLE = False

_missing = []
if plt is None:
    _missing.append("matplotlib")
if np is None:
    _missing.append("numpy")
if pd is None:
    _missing.append("pandas")
if _missing:
    missing_libs = ", ".join(_missing)
    raise RuntimeError(
        "eln_analysis.py requires the following third-party packages: "
        f"{missing_libs}. Install them (e.g. pip install matplotlib numpy pandas) or "
        "use src/eln_core.py for a dependency-light run."
    )

if plt is not None:
    plt.switch_backend("Agg")  # ensures headless rendering


@dataclass
class StockAssumption:
    ticker: str
    price: float
    volatility: float  # annualised, in decimal (e.g. 0.25)
    dividend_yield: float  # continuous yield proxy
    china_revenue_share: float  # for narrative context, not used in pricing
    commentary: str


@dataclass
class BasketSelection:
    first: str
    second: str
    weights: Tuple[float, float] = (0.5, 0.5)


@dataclass
class MarketParams:
    risk_free_rate: float  # continuously compounded, annualised
    time_to_maturity: float  # in years
    correlation: float  # pairwise correlation between the two equities


@dataclass
class PricingOutputs:
    note_price: float
    option_breakdown: Dict[str, float]
    basket_spot: float
    basket_vol: float
    basket_dividend_yield: float
    payoff_grid: 'pd.DataFrame'
    delta_profile_terminal: 'pd.DataFrame'
    delta_profile_t0: 'pd.DataFrame'


NOTIONAL = 10_000_000  # client capital


def black_scholes_call(
    spot: float,
    strike: float,
    rate: float,
    dividend_yield: float,
    volatility: float,
    time_to_maturity: float,
) -> Tuple[float, float]:
    """Return call price and delta under the Black-Scholes model.

    All rates are continuous and volatility is annualised. Delta is with respect
    to the underlying spot.
    """
    if spot <= 0 or strike <= 0:
        raise ValueError("Spot and strike must be positive.")
    if volatility <= 0:
        raise ValueError("Volatility must be positive.")
    sqrt_t = math.sqrt(time_to_maturity)
    d1 = (
        math.log(spot / strike)
        + (rate - dividend_yield + 0.5 * volatility ** 2) * time_to_maturity
    ) / (volatility * sqrt_t)
    d2 = d1 - volatility * sqrt_t
    nd1 = 0.5 * (1.0 + math.erf(d1 / math.sqrt(2.0)))
    nd2 = 0.5 * (1.0 + math.erf(d2 / math.sqrt(2.0)))
    discount_factor = math.exp(-rate * time_to_maturity)
    dividend_discount = math.exp(-dividend_yield * time_to_maturity)
    price = spot * dividend_discount * nd1 - strike * discount_factor * nd2
    delta = dividend_discount * nd1
    return price, delta


def build_stock_assumptions() -> Dict[str, StockAssumption]:
    """Assemble the four filter-selected equities with narrative inputs."""
    return {
        "LMT": StockAssumption(
            ticker="LMT",
            price=470.0,
            volatility=0.22,
            dividend_yield=0.027,
            china_revenue_share=0.03,
            commentary=(
                "Defense prime with minimal China revenue exposure; poised to benefit from "
                "elevated Indo-Pacific security spending under aggressive trade stances."
            ),
        ),
        "INTC": StockAssumption(
            ticker="INTC",
            price=34.0,
            volatility=0.35,
            dividend_yield=0.015,
            china_revenue_share=0.22,
            commentary=(
                "U.S.-based semiconductor manufacturer supported by CHIPS Act incentives "
                "and strategic realignment of advanced node supply chains."
            ),
        ),
        "CSCO": StockAssumption(
            ticker="CSCO",
            price=48.0,
            volatility=0.24,
            dividend_yield=0.029,
            china_revenue_share=0.05,
            commentary=(
                "Critical networking supplier benefiting from federal security procurement "
                "and restrictions on Chinese telecom equipment."
            ),
        ),
        "NUE": StockAssumption(
            ticker="NUE",
            price=170.0,
            volatility=0.30,
            dividend_yield=0.015,
            china_revenue_share=0.02,
            commentary=(
                "Domestic steel leader with electric arc capacity that benefits from U.S. "
                "infrastructure and reshoring themes amid tariff barriers."
            ),
        ),
    }


def choose_basket() -> BasketSelection:
    """Return the preferred two-name basket for the ELN."""
    return BasketSelection(first="LMT", second="INTC")


def effective_basket_inputs(
    selection: BasketSelection,
    stocks: Dict[str, StockAssumption],
    market: MarketParams,
) -> Tuple[float, float, float]:
    """Compute basket spot, dividend yield and volatility for a two-name basket."""
    w1, w2 = selection.weights
    s1 = stocks[selection.first]
    s2 = stocks[selection.second]
    basket_spot = w1 * s1.price + w2 * s2.price
    basket_dividend = (w1 * s1.price * s1.dividend_yield + w2 * s2.price * s2.dividend_yield) / basket_spot
    basket_variance = (
        (w1 * s1.volatility) ** 2
        + (w2 * s2.volatility) ** 2
        + 2 * w1 * w2 * s1.volatility * s2.volatility * market.correlation
    )
    basket_volatility = math.sqrt(basket_variance)
    return basket_spot, basket_dividend, basket_volatility


def price_note(
    stocks: Dict[str, StockAssumption],
    selection: BasketSelection,
    market: MarketParams,
    notional: float = NOTIONAL,
) -> PricingOutputs:
    """Price the ELN and assemble diagnostic artefacts."""
    basket_spot, basket_dividend, basket_volatility = effective_basket_inputs(selection, stocks, market)
    maturity = market.time_to_maturity

    discount_factor = math.exp(-market.risk_free_rate * maturity)
    weight_spot = notional / basket_spot

    # Option strikes expressed on the basket level
    strike_80 = 0.8 * basket_spot
    strike_100 = basket_spot

    call_80_price, call_80_delta = black_scholes_call(
        basket_spot,
        strike_80,
        market.risk_free_rate,
        basket_dividend,
        basket_volatility,
        maturity,
    )
    call_100_price, call_100_delta = black_scholes_call(
        basket_spot,
        strike_100,
        market.risk_free_rate,
        basket_dividend,
        basket_volatility,
        maturity,
    )

    option_80_cost = weight_spot * call_80_price
    option_100_cost = weight_spot * call_100_price

    note_price = (
        weight_spot * basket_spot  # intrinsic basket exposure
        + 0.2 * notional * discount_factor  # ZCB for 20% cushion
        - option_80_cost
        + 2.0 * option_100_cost
    )

    payoff_grid = build_payoff_grid(basket_spot)
    delta_profile_terminal = build_terminal_delta_profile()
    delta_profile_t0 = build_t0_delta_profile(
        basket_spot,
        basket_dividend,
        basket_volatility,
        market,
        weight_spot,
    )

    option_breakdown = {
        "Long basket": weight_spot * basket_spot,
        "Long 20% ZCB": 0.2 * notional * discount_factor,
        "Short call K=80%": -option_80_cost,
        "Long 2 calls K=100%": 2.0 * option_100_cost,
    }

    return PricingOutputs(
        note_price=note_price,
        option_breakdown=option_breakdown,
        basket_spot=basket_spot,
        basket_vol=basket_volatility,
        basket_dividend_yield=basket_dividend,
        payoff_grid=payoff_grid,
        delta_profile_terminal=delta_profile_terminal,
        delta_profile_t0=delta_profile_t0,
    )


def build_payoff_grid(basket_spot: float) -> pd.DataFrame:
    ratios = np.linspace(0.4, 1.6, 121)
    payoff = []
    for r in ratios:
        value = payoff_function(r)
        payoff.append(value)
    df = pd.DataFrame({
        "Basket_Performance": ratios,
        "Payoff_per_$1": payoff,
        "Payoff_per_$Notional": np.array(payoff) * NOTIONAL,
        "Basket_Level": ratios * basket_spot,
    })
    return df


def payoff_function(performance_ratio: float) -> float:
    if performance_ratio >= 1.0:
        return 2.0 * performance_ratio - 1.0
    if performance_ratio >= 0.8:
        return 1.0
    return performance_ratio + 0.2


def build_terminal_delta_profile() -> pd.DataFrame:
    ratios = np.linspace(0.4, 1.6, 121)
    terminal_delta = []
    for r in ratios:
        if r < 0.8:
            terminal_delta.append(1.0)
        elif r < 1.0:
            terminal_delta.append(0.0)
        else:
            terminal_delta.append(2.0)
    return pd.DataFrame({
        "Basket_Performance": ratios,
        "Terminal_Delta": terminal_delta,
    })


def build_t0_delta_profile(
    basket_spot: float,
    basket_dividend: float,
    basket_vol: float,
    market: MarketParams,
    weight_spot: float,
) -> pd.DataFrame:
    ratios = np.linspace(0.6, 1.4, 81)
    deltas = []
    prices = []
    for r in ratios:
        s = basket_spot * r
        call_80_price, call_80_delta = black_scholes_call(
            s,
            0.8 * basket_spot,
            market.risk_free_rate,
            basket_dividend,
            basket_vol,
            market.time_to_maturity,
        )
        call_100_price, call_100_delta = black_scholes_call(
            s,
            basket_spot,
            market.risk_free_rate,
            basket_dividend,
            basket_vol,
            market.time_to_maturity,
        )
        discount_factor = math.exp(-market.risk_free_rate * market.time_to_maturity)
        price = (
            weight_spot * s
            + 0.2 * NOTIONAL * discount_factor
            - weight_spot * call_80_price
            + 2.0 * weight_spot * call_100_price
        )
        delta = (
            weight_spot
            - weight_spot * call_80_delta
            + 2.0 * weight_spot * call_100_delta
        )
        prices.append(price)
        deltas.append(delta)
    return pd.DataFrame({
        "Basket_Ratio": ratios,
        "Note_Value": prices,
        "Delta_wrt_Basket": deltas,
    })


def export_excel(output: PricingOutputs, stocks: Dict[str, StockAssumption], selection: BasketSelection, market: MarketParams, path: str) -> None:
    writer = pd.ExcelWriter(path, engine="openpyxl")

    assumptions_df = pd.DataFrame([
        {
            "Ticker": s.ticker,
            "Spot": s.price,
            "Volatility": s.volatility,
            "Dividend_Yield": s.dividend_yield,
            "China_Revenue_Share": s.china_revenue_share,
            "Commentary": s.commentary,
        }
        for s in stocks.values()
    ])
    assumptions_df.to_excel(writer, sheet_name="Stock_Assumptions", index=False)

    pd.DataFrame([
        {
            "Parameter": "Risk-free rate",
            "Value": market.risk_free_rate,
        },
        {
            "Parameter": "Maturity (years)",
            "Value": market.time_to_maturity,
        },
        {
            "Parameter": "Correlation",
            "Value": market.correlation,
        },
        {
            "Parameter": "Basket spot",
            "Value": output.basket_spot,
        },
        {
            "Parameter": "Basket dividend yield",
            "Value": output.basket_dividend_yield,
        },
        {
            "Parameter": "Basket volatility",
            "Value": output.basket_vol,
        },
        {
            "Parameter": "Notional",
            "Value": NOTIONAL,
        },
        {
            "Parameter": "Selected tickers",
            "Value": f"{selection.first}, {selection.second}",
        },
        {
            "Parameter": "Note price",
            "Value": output.note_price,
        },
    ]).to_excel(writer, sheet_name="Market_Inputs", index=False)

    breakdown_df = pd.DataFrame(
        [(k, v) for k, v in output.option_breakdown.items()],
        columns=["Component", "Present_Value"],
    )
    breakdown_df.to_excel(writer, sheet_name="Pricing_Summary", index=False)

    output.payoff_grid.to_excel(writer, sheet_name="Payoff_Profile", index=False)
    output.delta_profile_terminal.to_excel(writer, sheet_name="Terminal_Delta", index=False)
    output.delta_profile_t0.to_excel(writer, sheet_name="Delta_T0", index=False)

    writer.close()


def export_charts(output: PricingOutputs, base_path: str | Path) -> Dict[str, str]:
    base_path = Path(base_path)
    base_path.mkdir(parents=True, exist_ok=True)
    payoff_path = base_path / "payoff_profile.png"
    delta_terminal_path = base_path / "terminal_delta.png"
    delta_t0_path = base_path / "delta_t0.png"

    plt.figure(figsize=(7, 4))
    plt.plot(output.payoff_grid["Basket_Performance"], output.payoff_grid["Payoff_per_$1"], label="ELN payoff")
    plt.axvline(0.8, color="grey", linestyle="--", linewidth=1)
    plt.axvline(1.0, color="grey", linestyle="--", linewidth=1)
    plt.xlabel("Basket performance (Final / Initial)")
    plt.ylabel("Payoff per $1 notional")
    plt.title("Bullish ELN Payoff")
    plt.legend()
    plt.grid(True, alpha=0.3)
    plt.tight_layout()
    plt.savefig(payoff_path, dpi=200)
    plt.close()

    plt.figure(figsize=(7, 4))
    plt.step(
        output.delta_profile_terminal["Basket_Performance"],
        output.delta_profile_terminal["Terminal_Delta"],
        where="post",
        label="Terminal delta",
    )
    plt.xlabel("Basket performance (Final / Initial)")
    plt.ylabel("dPayoff/dBasket")
    plt.title("Terminal Delta Profile")
    plt.grid(True, alpha=0.3)
    plt.legend()
    plt.tight_layout()
    plt.savefig(delta_terminal_path, dpi=200)
    plt.close()

    plt.figure(figsize=(7, 4))
    plt.plot(
        output.delta_profile_t0["Basket_Ratio"],
        output.delta_profile_t0["Delta_wrt_Basket"],
        label="Time-0 delta",
    )
    plt.xlabel("Current basket level / Initial")
    plt.ylabel("Delta vs basket")
    plt.title("Time-0 Delta Sensitivity")
    plt.grid(True, alpha=0.3)
    plt.legend()
    plt.tight_layout()
    plt.savefig(delta_t0_path, dpi=200)
    plt.close()

    return {
        "payoff": str(payoff_path),
        "terminal_delta": str(delta_terminal_path),
        "delta_t0": str(delta_t0_path),
    }


def export_ppt(output: PricingOutputs, stocks: Dict[str, StockAssumption], selection: BasketSelection, market: MarketParams, chart_paths: Dict[str, str], ppt_path: str) -> None:
    if not PPTX_AVAILABLE:
        raise RuntimeError("python-pptx is not available. Install it to build the presentation.")

    pres = Presentation()

    title_slide = pres.slides.add_slide(pres.slide_layouts[0])
    title_slide.shapes.title.text = "US-China Trade War Bullish ELN"
    title_slide.placeholders[1].text = "Structuring overview prepared for client"

    slide2 = pres.slides.add_slide(pres.slide_layouts[1])
    slide2.shapes.title.text = "Equity Filters"
    body = slide2.shapes.placeholders[1].text_frame
    body.text = (
        "Election outcome with elevated trade friction -> prioritise defence, "
        "reshoring, and secure tech supply chains. Filters: \n"
    )
    for stock in stocks.values():
        p = body.add_paragraph()
        p.text = f"{stock.ticker}: {stock.commentary}"
        p.level = 1

    slide3 = pres.slides.add_slide(pres.slide_layouts[5])
    slide3.shapes.title.text = "Payoff Mechanics"
    left = Inches(0.5)
    top = Inches(1.5)
    height = Inches(3.5)
    slide3.shapes.add_picture(chart_paths["payoff"], left, top, height=height)

    slide4 = pres.slides.add_slide(pres.slide_layouts[5])
    slide4.shapes.title.text = "Delta Profiles"
    height = Inches(3.0)
    slide4.shapes.add_picture(chart_paths["terminal_delta"], Inches(0.5), Inches(1.5), height=height)
    slide4.shapes.add_picture(chart_paths["delta_t0"], Inches(5.0), Inches(1.5), height=height)

    pres.save(ppt_path)


def main() -> None:

    generate_outputs()


def generate_outputs(output_dir: str | Path | None = None) -> PricingOutputs:
    """Generate calculation workbook, charts, and PPT to the given directory."""

    base_dir = Path(output_dir) if output_dir is not None else Path(__file__).resolve().parent.parent / "outputs"
    base_dir.mkdir(parents=True, exist_ok=True)

    stocks = build_stock_assumptions()
    selection = choose_basket()
    market = MarketParams(risk_free_rate=0.045, time_to_maturity=1.0, correlation=0.35)

    outputs = price_note(stocks, selection, market)

    excel_path = base_dir / "eln_calculations.xlsx"
    chart_paths = export_charts(outputs, base_dir)
    export_excel(outputs, stocks, selection, market, str(excel_path))

    ppt_path = base_dir / "eln_presentation.pptx"
    try:
        export_ppt(outputs, stocks, selection, market, chart_paths, str(ppt_path))
    except RuntimeError as exc:
        print(f"Warning: {exc}")

    summary = {
        "Note price": outputs.note_price,
        "Basket spot": outputs.basket_spot,
        "Basket vol": outputs.basket_vol,
        "Basket dividend": outputs.basket_dividend_yield,
    }
    print("Key results:")
    for k, v in summary.items():
        print(f"  {k}: {v:,.4f}")
    print("Component PVs:")
    for k, v in outputs.option_breakdown.items():
        print(f"  {k}: {v:,.2f}")

    return outputs


if __name__ == "__main__":
    main()
